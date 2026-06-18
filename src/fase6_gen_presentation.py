# -*- coding: utf-8 -*-
"""
FASE 6 — Generador de la presentación final ELI556 Tarea 2 (rediseño v2).

Narrativa centrada en el TRATAMIENTO DE DATOS y el FLUJO DE LOS CÓDIGOS:
- Lámina dedicada a la base de datos original NREL Cocoa y a la anatomía del CSV.
- Diagramas de flujo reales por fase (generados por scratch/gen_diagramas_flujo.py).
- Contenido alineado 1:1 con lo que el código ejecuta (IAM físico + factor
  espectral aplicados, albedo 0.20, ventana de 12 meses sin doble cobertura,
  energía integrada a 5 min, track de recurso real PVGIS, ajuste por
  minimize+bounds, α_Isc por fallback, etc.).
- Numeración de láminas y referencias cruzadas a anexos calculadas automáticamente.

Salida: output/Presentacion_Final_ELI556_Atacama_POA_Tarea2_revisado.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import os
import re

# ============================================================
# SISTEMA DE DISEÑO
# ============================================================
DARK_BG = RGBColor(0x0B, 0x0C, 0x10)
PANEL_BG = RGBColor(0x1F, 0x28, 0x33)
PANEL_BG2 = RGBColor(0x27, 0x31, 0x3F)
ACCENT_GOLD = RGBColor(0xF1, 0xC4, 0x0F)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x5E, 0x3A)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
TEXT_WHITE = RGBColor(0xFB, 0xFC, 0xFD)
TEXT_GREY = RGBColor(0x9A, 0xA5, 0xB1)
CODE_COLOR = RGBColor(0x7F, 0xDB, 0xFF)

MARGIN = Inches(0.4)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
GAP = Inches(0.3)
COL_W = (SLIDE_W - 2 * MARGIN - GAP) / 2
CONTENT_Y = Inches(1.30)

IMG = {
    'pipeline':   'output/Extra_Resultados/diagramas/pipeline_general.png',
    'estructura': 'output/Extra_Resultados/diagramas/estructura_csv.png',
    'embudo':     'output/Extra_Resultados/diagramas/embudo_datos.png',
    'ingesta_limpieza': 'output/Extra_Resultados/diagramas/ingesta_limpieza.png',
    'f1':         'output/Extra_Resultados/diagramas/flujo_fase1.png',
    'f2':         'output/Extra_Resultados/diagramas/flujo_fase2.png',
    'f3':         'output/Extra_Resultados/diagramas/flujo_fase3.png',
    'f4':         'output/Extra_Resultados/diagramas/flujo_fase4.png',
    'geo':        'output/Extra_Resultados/geo_emulation_map.png',
    'poa_cmp22':  'output/Extra_Resultados/poa_irradiance_cmp22.png',
    'circuito':   'output/Extra_Resultados/circuito_equivalente_sdm.png',
    'curvas_iv':  'output/Extra_Resultados/curvas_iv_pv_src.png',
    'degradacion':'output/Extra_Resultados/degradacion_termica_scatter.png',
    'pr_comp':    'output/Extra_Resultados/pr_mensual_comparativo.png',
    'perfil_dia': 'output/Extra_Resultados/perfil_dia_tipico.png',
    'temp_hist':  'output/Fase1_Resultados/temp_hist_HIT.png',
    'val_scatter':'output/Fase2_Resultados/val_scatter_HIT.png',
    'optica':     'output/Extra_Resultados/modificadores_opticos.png',
    'pr_real':    'output/Extra_Resultados/pr_mensual_atacama_real.png',
    'yield_real': 'output/Extra_Resultados/yield_atacama_real.png',
}

# ============================================================
# PLAN DE LÁMINAS (orden => numeración y referencias automáticas)
# ============================================================
ORDER = [
    'portada', 'intro', 'justificacion', 'marco',
    'pipeline', 'datos', 'ingesta', 'embudo',
    'emu_concepto', 'emu_impl', 'poa_medida',
    'f2_flujo', 'optica', 'dia_tipico', 'f3_flujo', 'f3_result', 'f4_flujo',
    'validacion', 'rs_n', 'perdidas_pr', 'veredicto_eco', 'recurso_real',
    'cumplimiento', 'limitaciones', 'conclusiones', 'referencias',
    'anexo_tools',
    'anexo1', 'anexo2', 'anexo3', 'anexo4', 'anexo5', 'anexo6', 'anexo7', 'anexo8',
    'anexo9', 'anexo10', 'anexo11', 'anexo12', 'anexo13',
]

SECTION = {  # (etiqueta kicker, color de sección)
    'intro': ('INTRODUCCIÓN', ACCENT_BLUE), 'justificacion': ('JUSTIFICACIÓN', ACCENT_BLUE),
    'marco': ('MARCO TEÓRICO', ACCENT_BLUE), 'pipeline': ('PIPELINE', ACCENT_GOLD),
    'datos': ('DATOS', ACCENT_GOLD), 'ingesta': ('DATOS', ACCENT_GOLD),
    'embudo': ('DATOS', ACCENT_GOLD), 'emu_concepto': ('EMULACIÓN', ACCENT_ORANGE),
    'emu_impl': ('EMULACIÓN', ACCENT_ORANGE), 'poa_medida': ('EMULACIÓN', ACCENT_ORANGE),
    'f2_flujo': ('MODELADO', ACCENT_GOLD), 'optica': ('MODELADO', ACCENT_GOLD),
    'dia_tipico': ('MODELADO', ACCENT_GOLD),
    'f3_flujo': ('MODELADO', ACCENT_GOLD), 'f3_result': ('MODELADO', ACCENT_GOLD),
    'f4_flujo': ('MODELADO', ACCENT_GOLD),
    'validacion': ('RESULTADOS', ACCENT_GREEN), 'rs_n': ('RESULTADOS', ACCENT_GREEN),
    'perdidas_pr': ('RESULTADOS', ACCENT_GREEN), 'veredicto_eco': ('RESULTADOS', ACCENT_GREEN),
    'recurso_real': ('RESULTADOS', ACCENT_GREEN),
    'cumplimiento': ('SÍNTESIS', ACCENT_GREEN), 'limitaciones': ('CIERRE', ACCENT_ORANGE),
    'conclusiones': ('CIERRE', ACCENT_ORANGE), 'referencias': ('CIERRE', ACCENT_ORANGE),
}
REF = {k: i + 1 for i, k in enumerate(ORDER)}
TOTAL = len(ORDER)

def A(key):
    """Referencia textual 'Anexo X, Lám. N' calculada automáticamente."""
    rom = {'anexo1': 'I', 'anexo2': 'II', 'anexo3': 'III', 'anexo4': 'IV',
           'anexo5': 'V', 'anexo6': 'VI', 'anexo7': 'VII', 'anexo8': 'VIII',
           'anexo9': 'IX', 'anexo10': 'X', 'anexo11': 'XI', 'anexo12': 'XII',
           'anexo13': 'XIII'}
    return f"(Anexo {rom[key]}, Lám. {REF[key]})"

# ============================================================
# HELPERS DE CONSTRUCCIÓN
# ============================================================

def apply_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def _emit_subsup(p, text, font_size, color, bold=False):
    """Emite runs procesando [_sub_] y ^sup^ (también dentro de negritas)."""
    for t in re.split(r'(\[_[^_\]]+_\]|\^[^^]+\^)', text):
        if not t:
            continue
        run = p.add_run()
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold:
            run.font.bold = True
        if t.startswith('[_') and t.endswith('_]'):
            run.text = t[2:-2]
            run.font.subscript = True
            run.font.size = Pt(max(8, font_size - 4))
        elif t.startswith('^') and t.endswith('^'):
            run.text = t[1:-1]
            run.font.superscript = True
            run.font.size = Pt(max(8, font_size - 4))
        else:
            run.text = t


def fmt_runs(p, text, font_size, default_color):
    """Mini-markdown seguro: **bold**, `code`, [_sub_], ^sup^ (sin cursivas _)."""
    for tok in re.split(r'(\*\*[^*]+\*\*|`[^`]+`)', text):
        if not tok:
            continue
        if tok.startswith('**') and tok.endswith('**'):
            _emit_subsup(p, tok[2:-2], font_size, default_color, bold=True)
        elif tok.startswith('`') and tok.endswith('`'):
            run = p.add_run()
            run.text = tok[1:-1]
            run.font.name = 'Consolas'
            run.font.size = Pt(max(8, font_size - 1))
            run.font.color.rgb = CODE_COLOR
        else:
            _emit_subsup(p, tok, font_size, default_color)


def add_par(tf, text, size=14, color=TEXT_WHITE, first=False, align=PP_ALIGN.LEFT,
            space_after=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    fmt_runs(p, text, size, color)
    if space_after is not None:
        p.space_after = Pt(space_after)
    return p


CURRENT_KEY = [None]  # clave de la lámina en construcción (para kicker/sección)


def add_title(slide, text):
    key = CURRENT_KEY[0]
    kick, accent = SECTION.get(key, (None, ACCENT_ORANGE))
    if key and key.startswith('anexo'):
        kick, accent = 'ANEXO', RGBColor(0x8A, 0x96, 0xA8)
    box = slide.shapes.add_textbox(MARGIN, MARGIN - Inches(0.05),
                                   SLIDE_W - 2 * MARGIN - Inches(2.0), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = add_par(tf, text, size=25, color=ACCENT_GOLD, first=True)
    for run in p.runs:
        run.font.bold = True
    if kick:
        chip_w = Inches(1.85)
        chip = slide.shapes.add_shape(1, SLIDE_W - MARGIN - chip_w, MARGIN + Inches(0.06),
                                      chip_w, Inches(0.34))
        chip.fill.solid(); chip.fill.fore_color.rgb = PANEL_BG2
        chip.line.color.rgb = accent; chip.line.width = Pt(1.2)
        ctf = chip.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        run = cp.add_run(); run.text = kick
        run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = accent
    line = slide.shapes.add_shape(1, MARGIN, MARGIN + Inches(0.72),
                                  SLIDE_W - 2 * MARGIN, Inches(0.022))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()


def add_footer(slide, num):
    box = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.4),
                                   SLIDE_W - 2 * MARGIN, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"ELI556 | Grupo Alta Tensión | Lámina {num}/{TOTAL}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.RIGHT


def add_panel(slide, left, top, width, height, title="", accent=None):
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL_BG
    rect.line.color.rgb = accent if accent else ACCENT_BLUE
    rect.line.width = Pt(1.0)
    if accent:
        bar = slide.shapes.add_shape(1, left, top, Inches(0.08), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    if title:
        box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.10),
                                       width - Inches(0.3), Inches(0.45))
        p = add_par(box.text_frame, title, size=18, color=ACCENT_GOLD, first=True)
        for run in p.runs:
            run.font.bold = True
    return rect


def _check_overflow(text, width, height, size, line_space):
    """Heurística: estima la altura del texto y avisa si excede la caja.
    No altera el render; solo imprime [OVERFLOW] durante la generación."""
    import math
    w_in = width / 914400
    h_in = height / 914400
    clean_total = 0
    lines = 0
    cpl = max(8, w_in * 72 / (size * 0.42))  # caracteres por línea estimados
    for par in text.split('\n'):
        clean = re.sub(r'\*\*|`|\[_|_\]|\^', '', par)
        lines += max(1, math.ceil(len(clean) / cpl))
        clean_total += len(clean)
    est_h = lines * size * 1.15 * (line_space or 1.0) / 72
    if est_h > h_in * 1.04:
        print(f"  [OVERFLOW] {CURRENT_KEY[0]}: texto de {clean_total} car. "
              f"~{est_h:.2f}in en caja de {h_in:.2f}in (size {size})")


def add_text(slide, left, top, width, height, text, size=14, color=TEXT_WHITE,
             align=PP_ALIGN.LEFT, line_space=None):
    _check_overflow(text, width, height, size, line_space)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = add_par(tf, line, size=size, color=color, first=(i == 0), align=align)
        if line_space:
            p.line_spacing = line_space
    return box


def panel_with_body(slide, left, top, width, height, title, body, accent=None,
                    body_size=17, line_space=1.10):
    add_panel(slide, left, top, width, height, title, accent)
    add_text(slide, left + Inches(0.22), top + Inches(0.70),
             width - Inches(0.44), height - Inches(0.88),
             body, size=body_size, line_space=line_space)


def stat_tile(slide, left, top, width, height, value, label, accent=ACCENT_GOLD):
    """Tarjeta de cifra destacada (jerarquía visual: número grande + etiqueta)."""
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid(); rect.fill.fore_color.rgb = PANEL_BG2
    rect.line.color.rgb = accent; rect.line.width = Pt(1.2)
    add_text(slide, left, top + Inches(0.08), width, Inches(0.75),
             "**" + value + "**", size=42, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.06), top + height - Inches(0.72),
             width - Inches(0.12), Inches(0.64),
             label, size=17.0, color=TEXT_GREY, align=PP_ALIGN.CENTER)


def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    # marcador de imagen faltante
    ph = slide.shapes.add_shape(1, left, top, width or Inches(3), height or Inches(2))
    ph.fill.solid(); ph.fill.fore_color.rgb = PANEL_BG2
    add_text(slide, left, top, width or Inches(3), Inches(0.5),
             f"[imagen no encontrada: {path}]", size=10, color=ACCENT_ORANGE)
    return None


def full_width_image(slide, path, max_w=12.5, top_in=1.35, max_bottom=6.95):
    """Inserta imagen centrada a lo ancho respetando el área útil."""
    from PIL import Image as PILImage
    if not os.path.exists(path):
        return add_image(slide, path, Inches(0.5), Inches(top_in))
    w_px, h_px = PILImage.open(path).size
    ar = w_px / h_px
    w = max_w
    h = w / ar
    if top_in + h > max_bottom:
        h = max_bottom - top_in
        w = h * ar
    left = (13.333 - w) / 2
    return slide.shapes.add_picture(path, Inches(left), Inches(top_in),
                                    width=Inches(w), height=Inches(h))


def add_caption(slide, text, top_in, h_in=0.85, accent=ACCENT_GOLD, size=17.0):
    left = MARGIN
    w = SLIDE_W - 2 * MARGIN
    rect = slide.shapes.add_shape(1, left, Inches(top_in), w, Inches(h_in))
    rect.fill.solid(); rect.fill.fore_color.rgb = PANEL_BG
    rect.line.color.rgb = accent; rect.line.width = Pt(1.0)
    # Use zero margin text frame to prevent text overflow at the bottom
    _check_overflow(text, w - Inches(0.4), Inches(h_in - 0.08), size, 1.08)
    box = slide.shapes.add_textbox(left + Inches(0.2), Inches(top_in + 0.04),
                                   w - Inches(0.4), Inches(h_in - 0.08))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, line in enumerate(text.split('\n')):
        p = add_par(tf, line, size=size, color=TEXT_WHITE, first=(i == 0))
        p.line_spacing = 1.08


def add_table(slide, left, top, width, height, headers, rows, col_ratios=None,
              fs_head=13, fs_body=12):
    n_rows, n_cols = len(rows) + 1, len(headers)
    gframe = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gframe.table
    if col_ratios:
        total = sum(col_ratios)
        for j, r in enumerate(col_ratios):
            table.columns[j].width = Emu(int(width * r / total))
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL_BG2
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        fmt_runs(p, htxt, fs_head, ACCENT_GOLD)
        for run in p.runs:
            run.font.bold = True
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL_BG if i % 2 == 0 else RGBColor(0x1A, 0x22, 0x2C)
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            fmt_runs(p, str(val), fs_body, TEXT_WHITE)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    return table


def add_eq(slide, formula, left, top, height, color='#E8A838', max_width=Inches(5.8)):
    """Renderiza una ecuación con mathtext de matplotlib (transparente)."""
    import matplotlib.pyplot as plt
    from PIL import Image as PILImage
    import hashlib
    tmp = 'output/temp_latex'
    os.makedirs(tmp, exist_ok=True)
    h = hashlib.md5(f"{formula}_{color}".encode()).hexdigest()
    fp = os.path.join(tmp, f"eq_{h}.png")
    if not os.path.exists(fp):
        fig = plt.figure(figsize=(20, 1.8), facecolor='none')
        ax = fig.add_axes([0, 0, 1, 1]); ax.axis('off')
        mt = formula if formula.startswith('$') else f"${formula}$"
        ax.text(0.5, 0.5, mt, fontsize=24, color=color, ha='center', va='center',
                math_fontfamily='cm')
        plt.savefig(fp, dpi=300, bbox_inches='tight', pad_inches=0.05, transparent=True)
        plt.close(fig)
        img = PILImage.open(fp)
        bbox = img.getbbox()
        if bbox:
            img.crop(bbox).save(fp)
    img = PILImage.open(fp)
    ar = img.width / img.height
    w = height * ar
    if w > max_width:
        w = max_width
        height = int(w / ar)
    slide.shapes.add_picture(fp, left, top, width=int(w), height=int(height))


# ============================================================
# CONSTRUCTORES DE LÁMINAS
# ============================================================

def s_portada(slide, n):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_ORANGE; bar.line.fill.background()
    add_text(slide, Inches(0.85), Inches(1.3), Inches(11.6), Inches(1.7),
             "**Evaluación de Tecnologías Fotovoltaicas en el Desierto de Atacama**",
             size=40, color=ACCENT_GOLD)
    add_text(slide, Inches(0.85), Inches(2.62), Inches(11.6), Inches(0.5),
             "Tratamiento de datos experimentales NREL y modelamiento dinámico m-Si vs HIT — Año 2026",
             size=17, color=TEXT_GREY)
    line = slide.shapes.add_shape(1, Inches(0.85), Inches(3.15), Inches(4.5), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT_BLUE; line.line.fill.background()
    panel = slide.shapes.add_shape(1, Inches(0.8), Inches(3.6), Inches(11.7), Inches(2.75))
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL_BG
    panel.line.color.rgb = ACCENT_BLUE; panel.line.width = Pt(1.5)
    add_text(slide, Inches(1.2), Inches(3.95), Inches(11.0), Inches(2.2),
             "Integrantes / Estudiantes:\n"
             "Laury Gualdron  |  Sebastian Marin  |  Alejandro Hernández\n"
             "\n"
             "Profesor Guía: Carlos Cardenas\n"
             "ELI556 — Modelamiento y Análisis de Sistemas PV  |  Grupo Alta Tensión (AT)\n"
             "Fecha de presentación: Jueves, 11 de junio de 2026", size=15)


def s_intro(slide, n):
    add_title(slide, "a) Introducción y contextualización: El Desierto de Atacama")
    H = Inches(5.45)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Contexto: El Recurso Solar más Extremo",
        "• **Irradiación excepcional:** GHI anual > 2,900 kWh/m² en el norte de "
        "Chile (máxima mundial) — motivación del sitio.\n\n"
        "• **Cielos limpios y altitud:** baja atenuación atmosférica a "
        "2,400 m.s.n.m.\n\n"
        "• **Sitio simulado:** San Pedro de Atacama "
        "(lat −22.91°, lon −68.20°, UTC−4).\n\n"
        "• **Base experimental:** mediciones NREL (Cocoa, Florida) emuladas "
        "estacionalmente al hemisferio sur (Limitaciones, Lám. %d)." % REF['limitaciones'],
        accent=ACCENT_BLUE, body_size=17.5)
    panel_with_body(slide, MARGIN + COL_W + GAP, CONTENT_Y, COL_W, H,
        "El Desafío: Estrés Térmico de Operación",
        "• **Calentamiento severo:** las celdas operan sobre 65 °C al mediodía "
        "solar de verano (máx. simulado: 73.3 °C en HIT).\n\n"
        "• **Degradación térmica:** la potencia máxima y el voltaje de circuito "
        "abierto caen al subir la temperatura de celda (T[_c_], máx. simulado 73.3 °C).\n\n"
        "• **Objetivo:** modelar registro a registro (cadencia 5 min) qué "
        "tecnología — **m-Si** o **HIT** — resiste mejor el estrés térmico "
        "desértico durante 2026, y traducirlo a PR anual.",
        accent=ACCENT_ORANGE, body_size=17.5)


def s_justificacion(slide, n):
    add_title(slide, "b) Justificación e impacto: Selección de Tecnologías")
    add_panel(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(5.7),
              "Familias disponibles en el Dataset Cocoa (NREL) y criterio de selección",
              accent=ACCENT_GOLD)
    headers = ["Tecnología", "Eficiencia", "Coef. Temp. (Pmp)", "Comportamiento en Desierto", "Decisión / Rol"]
    rows = [
        ["m-Si / x-Si\n(Silicio Monocristalino)", "17% - 21%", "−0.40 %/°C (Malo)",
         "Grandes pérdidas por calor (baja tolerancia).", "SELECCIONADO (Línea Base)"],
        ["HIT\n(Heterounión)", "20% - 22%", "−0.26 %/°C (Excelente)",
         "Mantiene alta producción bajo estrés térmico.", "SELECCIONADO (Premium)"],
        ["CdTe\n(Teluro de Cadmio)", "15% - 18%", "−0.28 %/°C (Excelente)",
         "Desempeño térmico sobresaliente; toxicidad por Cadmio.", "DESCARTADO (Menor contraste)"],
        ["CIGS\n(Película Delgada)", "14% - 16%", "−0.35 %/°C (Bueno)",
         "Pérdidas moderadas; susceptible a degradación por humedad.", "DESCARTADO (Sin contraste extremo)"],
        ["a-Si\n(Silicio Amorfo)", "6% - 10%", "−0.20 %/°C (Excelente)",
         "Muy baja eficiencia base y alta degradación inicial.", "DESCARTADO (Inviable comercial)"],
    ]
    add_table(slide, MARGIN + Inches(0.25), CONTENT_Y + Inches(0.55),
              SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(3.3),
              headers, rows, col_ratios=[2.2, 1.0, 1.4, 2.6, 1.9], fs_head=17, fs_body=15.0)
    add_text(slide, MARGIN + Inches(0.25), CONTENT_Y + Inches(4.55),
             SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.5),
             "Criterio: máximo contraste térmico entre tecnologías comerciales viables, con datos "
             "experimentales completos en el dataset (`mSi0166` y `HIT05667`).",
             size=15.5, color=TEXT_GREY)


def s_marco(slide, n):
    add_title(slide, "c) Marco referencial y revisión de literatura: Modelos Teóricos")
    H = Inches(5.45)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Literatura Académica y Modelos APLICADOS",
        "• **Modelo de un Diodo Simple (SDM) [De Soto et al., 2006]:** circuito "
        "equivalente con pérdidas óhmicas en serie (R[_s_]) y fugas shunt (R[_sh_]). "
        f"{A('anexo4')}\n\n"
        "• **Modelo Térmico de Sandia (SAPM) [King et al., 2004]:** estima la "
        f"temperatura interna de celda desde POA y viento. {A('anexo3')}\n\n"
        "• **Transposición de Perez (1990):** lleva GHI/DNI/DHI al plano inclinado "
        f"con cielo difuso anisotrópico. {A('anexo1')}\n\n"
        "• **Escalamiento De Soto (2006):** traslada los 5 parámetros desde STC a "
        f"condiciones operacionales vía bandgap. {A('anexo6')}\n\n"
        "• **Modificadores ópticos APLICADOS:** IAM físico de Snell-Bouguer y "
        f"factor espectral por masa de aire [King et al., 2004] — ejecutados en la "
        f"Fase 2 sobre la POA. {A('anexo2')}",
        accent=ACCENT_BLUE, body_size=15.0, line_space=1.05)
    
    left_panel_x = MARGIN + COL_W + GAP
    add_panel(slide, left_panel_x, CONTENT_Y, COL_W, H,
              "Circuito Equivalente de un Diodo (SDM)", accent=ACCENT_GOLD)
    
    # Enlarge the cropped circuit diagram image and center it horizontally in the panel
    img_w = Inches(5.9)
    img_h = Inches(2.60)  # Preserving the cropped aspect ratio (approx 2.27)
    add_image(slide, IMG['circuito'],
              left_panel_x + Inches(0.11), CONTENT_Y + Inches(0.65),
              width=img_w, height=img_h)
              
    # Add text block below the image to describe components and utilize space
    add_text(slide, left_panel_x + Inches(0.22), CONTENT_Y + Inches(3.35),
             COL_W - Inches(0.44), Inches(2.00),
             "Componentes del circuito equivalente (SDM):\n"
             "• **I[_L_] (Corriente fotogenerada):** proporcional a la irradiancia.\n"
             "• **Diodo (D):** representa la recombinación en la unión P-N.\n"
             "• **R[_sh_] (Resistencia shunt):** corriente de fuga interna.\n"
             "• **R[_s_] (Resistencia serie):** pérdidas óhmicas de contacto.\n"
             "• **I / V:** corriente y voltaje de salida entregados.",
             size=15.0, color=TEXT_WHITE, line_space=1.08)


def s_pipeline(slide, n):
    add_title(slide, "Procedimiento: Arquitectura del Pipeline de Simulación (Fases 0–4)")
    full_width_image(slide, IMG['pipeline'], max_w=12.5, top_in=1.30, max_bottom=5.50)
    add_caption(slide,
        "**Flujo lógico de la simulación (Fases 0-4):** El proceso inicia con la exploración y emulación "
        "estacional de los datos (Fases 0-1) para adaptarlos al hemisferio sur. Luego, se limpia y calcula "
        "el recurso solar efectivo (Fase 2) y se extraen los parámetros físicos de los módulos (Fase 3). "
        "Finalmente, se realiza la simulación horaria dinámica (Fase 4) para obtener las series temporales de potencia y PR.",
        top_in=5.60, h_in=1.30, accent=ACCENT_GOLD)


def s_datos(slide, n):
    add_title(slide, "Procedimiento: La Base de Datos Original — NREL Cocoa, Florida")
    TW = (SLIDE_W - 2 * MARGIN - 3 * Inches(0.22)) / 4
    tiles = [("11", "módulos PV medidos\nen simultáneo", ACCENT_BLUE),
             ("5 min", "cadencia de medición\n(solo horas de sol)", ACCENT_GOLD),
             ("~420 mil", "curvas I-V completas\n(ene-2011 → mar-2012)", ACCENT_ORANGE),
             ("~1.2 GB", "en 11 archivos CSV\n(uno por módulo)", ACCENT_GREEN)]
    for i, (v, l, c) in enumerate(tiles):
        stat_tile(slide, MARGIN + i * (TW + Inches(0.22)), CONTENT_Y,
                  TW, Inches(1.45), v, l, accent=c)
    Y2 = CONTENT_Y + Inches(1.70)
    H2 = Inches(3.75)
    panel_with_body(slide, MARGIN, Y2, COL_W, H2,
        "¿Qué es?",
        "• Campaña del **NREL** para validar modelos PV [Marion et al., 2014] "
        "en **Cocoa, Florida** (28.39° N, 12 m.s.n.m.).\n\n"
        "• **Cada registro = una curva I-V completa** + meteorología "
        "instantánea + incertidumbres.\n\n"
        "• Seleccionados: `mSi0166` (**36,765 curvas**) y `HIT05667` "
        "(**38,377 curvas**) — criterio en Lám. %d." % REF['justificacion'],
        accent=ACCENT_BLUE, body_size=16.5)
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
        "¿Qué variables usamos?",
        "• **Usadas (11 columnas):** timestamp, Isc, Pmp, Imp, Vmp, Voc, "
        "T bulbo seco, presión, DNI, GHI, DHI.\n\n"
        "• **Verificación:** POA del piranómetro CMP22 clase A (Lám. %d).\n\n"
        "• **Descartadas:** cola de pares I-V crudos y columnas de QA — "
        "diccionario completo en %s." % (REF['poa_medida'], A('anexo8')),
        accent=ACCENT_ORANGE, body_size=16.5)


def s_anexo9(slide, n):
    add_title(slide, "Anexo IX: Anatomía del Archivo CSV y Decisiones de Ingesta")
    full_width_image(slide, IMG['estructura'], max_w=12.5, top_in=1.25, max_bottom=5.50)
    add_caption(slide,
        "**Estructura y lectura de datos:** Los archivos originales del NREL contienen una cabecera de metadatos "
        "y registros cada 5 minutos. Dado que cada fila incluye las curvas eléctricas crudas completas, su tamaño "
        "es variable y pesado (~110 MB). Para optimizar la memoria, el código lee el archivo en flujo (streaming) "
        "e indexa únicamente las 11 variables meteorológicas y eléctricas necesarias, descartando el resto al vuelo.",
        top_in=5.60, h_in=1.30, accent=ACCENT_BLUE)


def s_ingesta(slide, n):
    add_title(slide, "Procedimiento: Ingesta y Limpieza de la Base de Datos")
    full_width_image(slide, IMG['ingesta_limpieza'], max_w=12.5, top_in=1.20, max_bottom=5.20)
    add_caption(slide,
        "**Proceso de ingesta y depuración:** Cada fila del CSV mezcla 43 columnas fijas con una cola de pares I-V "
        "de largo variable y sentinelas `-9999` — anatomía detallada en %s —, lo que impide `pandas.read_csv` directo. "
        "El pipeline implementa un lector en flujo (streaming) que indexa solo las 11 variables útiles; luego "
        "reemplaza `-9999` por `NaN`, descarta filas nulas y recorta irradiancias negativas, conservando "
        "**95.8 %% (m-Si) y 96.1 %% (HIT)** de la ventana de 12 meses para el modelado." % A('anexo9'),
        top_in=5.30, h_in=1.60, accent=ACCENT_GOLD, size=16.5)


def s_embudo(slide, n):
    add_title(slide, "Procedimiento: Embudo de Datos — Trazabilidad por Fase")
    full_width_image(slide, IMG['embudo'], max_w=10.9, top_in=1.32, max_bottom=7.0)


def s_emu_concepto(slide, n):
    add_title(slide, "Procedimiento: Filtro de Emulación Geográfica — Concepto")
    full_width_image(slide, IMG['geo'], max_w=11.6, top_in=1.40, max_bottom=6.55)
    add_text(slide, MARGIN, Inches(6.65), SLIDE_W - 2 * MARGIN, Inches(0.5),
             "¿Por qué? Usar datos del hemisferio norte sin desfasar crearía solsticios "
             "invertidos: el filtro alinea las estaciones antes de simular.",
             size=14, color=TEXT_GREY, align=PP_ALIGN.CENTER)


def s_anexo10(slide, n):
    add_title(slide, "Anexo X: Florida (origen de los datos) vs Atacama (sitio simulado)")
    H = Inches(4.25)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Cocoa, Florida — clima original",
        "• Clima: subtropical húmedo (marítimo), 12 m.s.n.m.\n\n"
        "• Irradiancia difusa alta (~35–45 %) por humedad y nubosidad "
        "convectiva frecuente.\n\n"
        "• DNI máximo ~950 W/m² (atenuado).\n\n"
        "• T[_c_] pico de celda: ~55–60 °C.",
        accent=ACCENT_BLUE, body_size=17.5, line_space=1.15)
    panel_with_body(slide, MARGIN + COL_W + GAP, CONTENT_Y, COL_W, H,
        "San Pedro de Atacama — sitio emulado",
        "• Clima: hiperárido (desértico extremo), 2,400 m.s.n.m.\n\n"
        "• Irradiancia difusa muy baja (~10–15 %): cielos permanentemente "
        "limpios.\n\n"
        "• DNI máximo real > 1,250 W/m² (extremo).\n\n"
        "• T[_c_] pico simulada: ~65–73 °C (alto estrés).",
        accent=ACCENT_GOLD, body_size=17.5, line_space=1.15)
    add_caption(slide,
        "**Límite explícito del método:** el filtro alinea estaciones, geometría solar e "
        "inclinación óptima (tilt = latitud, azimut 0° N); las **magnitudes medidas** de "
        "irradiancia y clima siguen siendo las de Florida (ver Limitaciones, Lám. %d)." % REF['limitaciones'],
        top_in=5.60, h_in=1.45, accent=ACCENT_ORANGE)


def s_emu_impl(slide, n):
    add_title(slide, "Procedimiento: Filtro de Emulación — Implementación (Fase 1)")
    full_width_image(slide, IMG['f1'], max_w=11.0, top_in=1.30, max_bottom=4.90)
    add_caption(slide,
        "**Proceso de emulación temporal:** el código reescribe el sitio en los metadatos (lat −22.91°, "
        "2,400 m, tilt 22.9°, azimut 0° N) y aplica un desfase de +6 meses a las fechas, proyectando el año a "
        "2026 **sin alterar los valores físicos medidos**. Se conserva solo la **primera ventana de 12 meses** "
        "[ene-2011, ene-2012): así cada mes se cubre una vez y se elimina la doble cobertura de jul–sep "
        "— contraste climático en %s." % A('anexo10'),
        top_in=5.00, h_in=1.18, accent=ACCENT_GOLD, size=15.5)
    add_caption(slide,
        "**Límite explícito del método:** se alinean estaciones y geometría solar; las **magnitudes medidas** "
        "de irradiancia y clima siguen siendo las de Florida (ver Limitaciones, Lám. %d)." % REF['limitaciones'],
        top_in=6.28, h_in=0.78, accent=ACCENT_ORANGE, size=15.5)


def s_poa_medida(slide, n):
    add_title(slide, "Procedimiento: Verificación del Recurso Medido (POA CMP22)")
    TW = Inches(3.85)
    gap = Inches(0.49)
    stat_tile(slide, MARGIN, CONTENT_Y, TW, Inches(1.50),
              "421,664", "registros de irradiancia válidos\n(11 archivos Cocoa)", accent=ACCENT_BLUE)
    stat_tile(slide, MARGIN + TW + gap, CONTENT_Y, TW, Inches(1.50),
              "1,443 W/m²", "POA máxima medida en el\nplano del arreglo", accent=ACCENT_GOLD)
    stat_tile(slide, MARGIN + 2 * (TW + gap), CONTENT_Y, TW, Inches(1.50),
              "CMP22", "piranómetro de referencia\nclase A (alta precisión)", accent=ACCENT_ORANGE)
    full_width_image(slide, IMG['poa_cmp22'], max_w=11.4, top_in=2.95, max_bottom=6.55)
    add_text(slide, MARGIN, Inches(6.62), SLIDE_W - 2 * MARGIN, Inches(0.45),
             "Las 11 curvas se superponen → consistencia instrumental antes de emular y modelar "
             "— notas metodológicas en %s." % A('anexo11'),
             size=14, color=TEXT_GREY, align=PP_ALIGN.CENTER)


def s_anexo11(slide, n):
    add_title(slide, "Anexo XI: Verificación del Recurso — Notas Metodológicas")
    
    # Tarjetas de Cifras Destacadas (Stat Tiles)
    TW = Inches(3.85)
    gap = Inches(0.49)
    stat_tile(slide, MARGIN, CONTENT_Y, TW, Inches(1.50),
              "421,664", "Registros de irradiancia válidos\n(11 archivos Cocoa)", accent=ACCENT_BLUE)
    stat_tile(slide, MARGIN + TW + gap, CONTENT_Y, TW, Inches(1.50),
              "1,443 W/m²", "Irradiancia POA máxima medida\nen el plano del arreglo", accent=ACCENT_GOLD)
    stat_tile(slide, MARGIN + 2 * (TW + gap), CONTENT_Y, TW, Inches(1.50),
              "CMP22", "Piranómetro de referencia\nclase A (alta precisión)", accent=ACCENT_ORANGE)
    
    # Paneles de Notas
    Y2 = CONTENT_Y + Inches(1.68)
    H2 = Inches(3.92)
    panel_with_body(slide, MARGIN, Y2, COL_W, H2,
        "Análisis de Consistencia de Datos",
        "• **Superposición de mediciones:** Las curvas de irradiancia diaria se "
        "superponen con exactitud, confirmando la consistencia instrumental.\n\n"
        "• **Rango temporal:** Campaña desde el **21 de enero de 2011** al "
        "**4 de marzo de 2012**.\n\n"
        "• **Tratamiento:** Mediana diaria suavizada (ventana 7 días) por archivo "
        "para aislar ruidos y nubes transitorias.",
        accent=ACCENT_BLUE, body_size=15.0, line_space=1.12)
    
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
        "Propósito del Control de Calidad",
        "• **Validación previa:** Confirmar la homogeneidad física de la radiación "
        "**antes** de la emulación geográfica y el modelado.\n\n"
        "• **Soiling y mantenimiento:** La coincidencia de curvas demuestra que "
        "no hay desviaciones locales significativas en los sensores de referencia.\n\n"
        "• **Alineación de sensores:** Garantizar que la irradiancia incidente en "
        "los 11 planos sea idéntica para una comparación justa.",
        accent=ACCENT_GOLD, body_size=15.0, line_space=1.12)


def s_f2_flujo(slide, n):
    add_title(slide, "Procedimiento: Recurso Solar Sintético y Perfil Térmico (Fase 2)")
    full_width_image(slide, IMG['f2'], max_w=11.5, top_in=1.35, max_bottom=5.20)
    add_caption(slide,
        "**Estimación de la radiación y temperatura:** Tras filtrar y limpiar los registros nulos, "
        "el código transpone las irradiancias al plano inclinado (POA) con el modelo difuso de Perez "
        "(albedo **0.20**). Sobre esa POA se aplican los modificadores ópticos IAM y espectral (lámina "
        "siguiente). Luego, con el modelo de Sandia (SAPM) y viento fijo de 1 m/s (peor caso de ventilación), "
        "se estima la temperatura de celda de cada tecnología (m-Si y HIT).",
        top_in=5.30, h_in=1.70, accent=ACCENT_ORANGE)


def s_dia_tipico(slide, n):
    add_title(slide, "Resultado intermedio: Un Día Despejado de Verano (28-ene-2026)")
    full_width_image(slide, IMG['perfil_dia'], max_w=11.8, top_in=1.30, max_bottom=5.35)
    add_caption(slide,
        "**El puente entre el dato y la física:** la temperatura de celda sigue fielmente a la irradiancia: al "
        "mediodía la POA roza los **1,000 W/m²** y la celda supera los **60 °C** — 35 °C sobre la condición "
        "nominal STC. En los días extremos del año la celda llega a **70–73 °C** y HIT acumula 430 registros "
        "sobre 65 °C — distribución en %s. Cada grado extra recorta la potencia útil; cuantificar ese recorte "
        "por tecnología es el objetivo de la sección Resultados." % A('anexo12'),
        top_in=5.45, h_in=1.55, accent=ACCENT_ORANGE, size=15.5)


def s_f3_flujo(slide, n):
    add_title(slide, "Procedimiento: Extracción de Parámetros De Soto (Fase 3)")
    full_width_image(slide, IMG['f3'], max_w=11.5, top_in=1.35, max_bottom=5.20)
    add_caption(slide,
        "**Caracterización física de los módulos:** Para modelar el comportamiento eléctrico, se estiman los coeficientes "
        "térmicos de corriente y voltaje. Dado que la corriente arrojó coeficientes negativos debido a la estacionalidad, "
        "se aplicó un valor estándar de resguardo de +0.05 %/°C. Luego, los datos se trasladan a condiciones nominales (STC) "
        "y un optimizador no lineal ajusta los 5 parámetros de De Soto (corrientes, resistencias y factor de idealidad).",
        top_in=5.30, h_in=1.70, accent=ACCENT_GOLD)


def s_f3_result(slide, n):
    add_title(slide, "Procedimiento: Parámetros de Referencia Extraídos en STC")
    W1 = Inches(6.1)
    add_panel(slide, MARGIN, CONTENT_Y, W1, Inches(5.5),
              "Parámetros nominales ajustados (SRC)", accent=ACCENT_GOLD)
    headers = ["Parámetro", "m-Si (mSi0166)", "HIT (HIT05667)"]
    rows = [
        ["I[_L,ref_] (A)", "2.770", "5.620"],
        ["I[_0,ref_] (A)", "4.22 × 10^−9^", "4.42 × 10^−10^"],
        ["a[_ref_] (V) → n[_I_]", "1.110 → 1.20", "2.220 → 1.20"],
        ["R[_s,ref_] (Ω)", "0.010 (≈ inicial)", "0.010 (≈ inicial)"],
        ["R[_sh,ref_] (Ω)", "1,000 (≈ inicial)", "1,000 (≈ inicial)"],
        ["N[_s_] (celdas serie)", "36", "72"],
        ["α[_Isc_] (%/°C)", "+0.05 (resguardo)", "+0.05 (resguardo)"],
        ["β[_Voc_] (%/°C)", "−0.30 (medido)", "−0.22 (medido)"],
        ["I[_sc,ref_] / V[_oc,ref_]", "2.77 A / 22.5 V", "5.62 A / 51.6 V"],
        ["P[_STC_] del SDM (W)", "50.12", "237.96"],
    ]
    add_table(slide, MARGIN + Inches(0.2), CONTENT_Y + Inches(0.62),
              W1 - Inches(0.4), Inches(4.7), headers, rows,
              col_ratios=[1.5, 1.1, 1.1], fs_head=14.0, fs_body=12.5)
    X2 = MARGIN + W1 + Inches(0.25)
    W2 = SLIDE_W - X2 - MARGIN
    add_panel(slide, X2, CONTENT_Y, W2, Inches(3.05),
              "Curvas I-V / P-V reconstruidas en SRC", accent=ACCENT_BLUE)
    add_image(slide, IMG['curvas_iv'], X2 + Inches(0.3), CONTENT_Y + Inches(0.60),
              width=Inches(5.65), height=Inches(2.30))
    panel_with_body(slide, X2, CONTENT_Y + Inches(3.25), W2, Inches(2.30),
        "Lectura física",
        "• I[_0_] de HIT es **un orden de magnitud menor** → menor recombinación "
        "→ mayor V[_oc_] y mejor tolerancia térmica.\n"
        "• R[_s_] y R[_sh_] quedan **anclados a su inicialización**: con solo 3 "
        "condiciones (Isc, Voc, MPP) el residuo apenas depende de ellos — "
        "evidencia del mal condicionamiento que se discute en Lám. %d." % REF['rs_n'],
        accent=ACCENT_GOLD, body_size=14.0, line_space=1.08)


def s_f4_flujo(slide, n):
    add_title(slide, "Procedimiento: Simulación Anual y Performance Ratio (Fase 4)")
    full_width_image(slide, IMG['f4'], max_w=11.5, top_in=1.35, max_bottom=5.20)
    add_caption(slide,
        "**Cálculo de la producción eléctrica anual:** Con los parámetros calibrados, se evalúa el modelo "
        "para cada uno de los registros del año 2026. A partir de la irradiancia y temperatura de celda de cada instante, "
        "se escala el circuito de un diodo y se resuelve analíticamente la potencia útil máxima (Pmp). Finalmente, "
        "la energía anual integrada se compara con el rendimiento ideal para obtener el Performance Ratio (PR).",
        top_in=5.30, h_in=1.70, accent=ACCENT_GOLD)


def s_anexo12(slide, n):
    add_title(slide, "Anexo XII: Distribución Anual de Temperatura de Celda")
    add_panel(slide, MARGIN, CONTENT_Y, Inches(6.9), Inches(5.5),
              "Distribución de Temperatura de Celda (HIT)", accent=ACCENT_BLUE)
    add_image(slide, IMG['temp_hist'], MARGIN + Inches(0.55), CONTENT_Y + Inches(0.68),
              width=Inches(5.9), height=Inches(4.41))
    X2 = MARGIN + Inches(7.1)
    W2 = SLIDE_W - X2 - MARGIN
    panel_with_body(slide, X2, CONTENT_Y, W2, Inches(1.55),
        "Picos Térmicos Extremos",
        "**65–73 °C de operación recurrente** al mediodía solar de verano "
        "(HIT: 430 registros sobre 65 °C).",
        accent=ACCENT_ORANGE, body_size=16.5)
    panel_with_body(slide, X2, CONTENT_Y + Inches(1.75), W2, Inches(3.75),
        "Análisis Crítico de Resultados",
        "• **Estrés térmico evidenciado:** operación sostenida muy por encima "
        "de los 25 °C de STC.\n\n"
        "• **Enfriamiento convectivo limitado:** viento de 1 m/s en SAPM — "
        "escenario de mayor estrés físico.\n\n"
        "• **Consecuencia en voltaje:** el calor eleva I[_0_] exponencialmente, "
        "deprimiendo V[_oc_] y la potencia útil [De Soto et al., 2006].",
        accent=ACCENT_GOLD, body_size=15.5, line_space=1.10)


def s_validacion(slide, n):
    add_title(slide, "d) Desarrollo y discusión: Validación del Modelo Eléctrico")
    add_panel(slide, MARGIN, CONTENT_Y, Inches(6.9), Inches(5.5),
              "P[_mp_] medida vs simulada (HIT, mismas condiciones)", accent=ACCENT_BLUE)
    add_image(slide, IMG['val_scatter'], MARGIN + Inches(0.55), CONTENT_Y + Inches(0.68),
              width=Inches(5.9), height=Inches(4.41))
    X2 = MARGIN + Inches(7.1)
    W2 = SLIDE_W - X2 - MARGIN
    panel_with_body(slide, X2, CONTENT_Y, W2, Inches(1.55),
        "Precisión del Modelo de 5 Parámetros",
        "**R² ≈ 0.99**  |  RMSE ≈ 3 % de P[_STC_] (1.4 W m-Si · 6.5 W HIT).",
        accent=ACCENT_BLUE, body_size=17.5)
    panel_with_body(slide, X2, CONTENT_Y + Inches(1.75), W2, Inches(3.75),
        "¿Qué valida exactamente este gráfico?",
        "• **Misma meteorología:** cada punto compara la P[_mp_] **medida** por "
        "el trazador NREL con la **simulada** por el SDM bajo el mismo G y "
        "T[_c_] — no es una comparación entre sitios.\n\n"
        "• **Fidelidad del circuito:** R² ≈ 0.99 confirma que los 5 parámetros "
        "reproducen el comportamiento eléctrico real.\n\n"
        "• **Discrepancia:** leve dispersión a baja irradiancia "
        "(idealización R[_sh_] ∝ 1/G).",
        accent=ACCENT_GOLD, body_size=15.0, line_space=1.08)


def s_rs_n(slide, n):
    add_title(slide, "d) Desarrollo: Puntos Conflictivos — Acoplamiento R[_s_] – n")
    H = Inches(5.45)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Acoplamiento R[_s_] – factor de idealidad (n)",
        "• **Alta correlación paramétrica [De Soto et al., 2006]:** R[_s_] y n "
        "influyen de manera similar en la curvatura de la I-V cerca del MPP.\n\n"
        "• **Problema mal condicionado:** múltiples pares (R[_s_], n) logran un "
        "residuo mínimo casi idéntico, perdiendo significado físico.\n\n"
        "• **Mitigación aplicada:** bounds físicos en el optimizador — "
        "R[_s_] ∈ [0.001, 2] Ω, n[_I_] ∈ [0.6, 2.4] — más inicialización "
        "analítica (a₀ = N[_s_]·1.2·kT/q).\n\n"
        "• **Evidencia en este estudio:** n[_I_] ≈ 1.20 en ambas tecnologías y "
        "R[_s_]/R[_sh_] anclados a su inicialización (Lám. %d): el ajuste con 3 "
        "condiciones no los identifica de forma única." % REF['f3_result'],
        accent=ACCENT_BLUE, body_size=15.5, line_space=1.08)
    panel_with_body(slide, MARGIN + COL_W + GAP, CONTENT_Y, COL_W, H,
        "Limitaciones y Suposiciones del Modelo",
        "• **R[_sh_] e irradiancia:** la relación R[_sh_] = R[_sh,ref_]·(S[_ref_]/S) "
        "es empírica; a muy baja irradiancia puede sobreestimar la resistencia. "
        "Impacto: despreciable en el PR anual.\n\n"
        "• **R[_s_] constante con la temperatura:** ignora el cambio de "
        "resistividad del semiconductor con T[_c_]. De Soto (2006) validó errores "
        "de potencia < 2 % frente a datos NIST.\n\n"
        "• **Bandgap lineal con la temperatura:** se asume dE[_g_]/dT constante "
        "(−0.0002677 eV/K), simplificando la relación real de Varshni. Error "
        "inducido < 1 % en el rango operacional (0–70 °C).",
        accent=ACCENT_GOLD, body_size=17.0, line_space=1.12)


def s_perdidas_pr(slide, n):
    add_title(slide, "d) Desarrollo: Pérdidas Térmicas y Degradación de Potencia")
    add_panel(slide, MARGIN, CONTENT_Y, Inches(7.4), Inches(5.5),
              "Potencia normalizada vs Temperatura de Celda", accent=ACCENT_BLUE)
    add_image(slide, IMG['degradacion'], MARGIN + Inches(0.35), CONTENT_Y + Inches(0.70),
              width=Inches(6.8), height=Inches(4.04))
    X2 = MARGIN + Inches(7.6)
    W2 = SLIDE_W - X2 - MARGIN
    panel_with_body(slide, X2, CONTENT_Y, W2, Inches(1.55),
        "Coeficientes de Temperatura (P[_mp_])",
        "m-Si: **−0.40 %/°C** (pérdidas severas)\nHIT: **−0.26 %/°C** (alta tolerancia)",
        accent=ACCENT_ORANGE, body_size=17.5)
    panel_with_body(slide, X2, CONTENT_Y + Inches(1.75), W2, Inches(3.75),
        "Observación Científica",
        "• **HIT (naranja):** pendiente térmica suave — preserva alta conversión "
        "incluso a 65 °C.\n\n"
        "• **m-Si (azul):** descenso pronunciado que penaliza la potencia justo "
        "al mediodía solar (máxima irradiancia).\n\n"
        "• **Consecuencia:** esta divergencia térmica explica mecánicamente la "
        "brecha de PR anual entre tecnologías %s." % A('anexo7'),
        accent=ACCENT_GOLD, body_size=16.0, line_space=1.10)


def s_veredicto_eco(slide, n):
    add_title(slide, "d) Veredicto Técnico: Performance Ratio y Lectura Económica")
    TW = (COL_W - Inches(0.4)) / 2
    stat_tile(slide, MARGIN, CONTENT_Y, TW, Inches(1.5),
              "84.18 %", "PR anual HIT\n(mensual 81.6–88.2 %)", accent=ACCENT_GOLD)
    stat_tile(slide, MARGIN + TW + Inches(0.4), CONTENT_Y, TW, Inches(1.5),
              "81.61 %", "PR anual m-Si\n(mensual 78.7–86.9 %)", accent=ACCENT_BLUE)
    stat_tile(slide, MARGIN, CONTENT_Y + Inches(1.70), COL_W, Inches(1.50),
              "+2.57 puntos", "ventaja neta de HIT en PR (track validado, R² 0.99)", accent=ACCENT_GREEN)
    panel_with_body(slide, MARGIN, CONTENT_Y + Inches(3.35), COL_W, Inches(2.10),
        "Lectura técnica (recurso emulado, validado)",
        "• Menor coeficiente térmico de HIT (−0.26 vs −0.40 %%/°C): protege la "
        "potencia al mediodía.\n"
        "• El PR ya penaliza la pérdida óptica (≈3 %%) y espectral.\n"
        "• Economía con recurso real (Lám. %d): **+3,294 MWh, +USD 148k/año** "
        "en 100 MWp." % REF['recurso_real'],
        accent=ACCENT_GOLD, body_size=13.5, line_space=1.04)
    X2 = MARGIN + COL_W + GAP
    add_panel(slide, X2, CONTENT_Y, COL_W, Inches(5.45),
              "PR mensual comparativo (2026)", accent=ACCENT_BLUE)
    add_image(slide, IMG['pr_comp'], X2 + Inches(0.22), CONTENT_Y + Inches(1.45),
              width=Inches(5.85), height=Inches(2.80))


def s_anexo13(slide, n):
    add_title(slide, "Anexo XIII: Lectura Económica con Recurso REAL (Planta de 100 MWp)")
    TW = (SLIDE_W - 2 * MARGIN - 2 * Inches(0.3)) / 3
    tiles = [("+3,294 MWh", "energía adicional anual de HIT\n(recurso real Atacama, 100 MWp)", ACCENT_GOLD),
             ("+USD 148k", "facturación neta anual extra\n(a 45 USD/MWh)", ACCENT_GREEN),
             ("33 kWh/kWp", "diferencia de yield específico\nHIT − m-Si (PVGIS TMY)", ACCENT_BLUE)]
    for i, (v, l, c) in enumerate(tiles):
        stat_tile(slide, MARGIN + i * (TW + Inches(0.3)), CONTENT_Y,
                  TW, Inches(1.5), v, l, accent=c)
    Y2 = CONTENT_Y + Inches(1.75)
    H2 = Inches(3.70)
    panel_with_body(slide, MARGIN, Y2, COL_W, H2,
        "Supuestos explícitos del cálculo",
        "• **Recurso REAL (no escalado):** PVGIS TMY de San Pedro de Atacama, "
        "POA ≈ 2,810 kWh/m²·año (track de la Lám. %d).\n\n"
        "• ΔPR +1.17 pts (HIT 84.99 %% vs m-Si 83.82 %%) → Δyield 33 kWh/kWp·año.\n\n"
        "• 33 kWh/kWp × 100,000 kWp = **+3,294 MWh/año**.\n\n"
        "• Precio de venta ≈ 45 USD/MWh → **+USD 148k/año**." % REF['recurso_real'],
        accent=ACCENT_GOLD, body_size=16.5)
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
        "Conclusión de diseño",
        "• La prima CAPEX de los módulos HIT se **amortiza** por su mayor yield "
        "en el recurso desértico de clase mundial (≈2,370 kWh/kWp).\n\n"
        "• La ventaja relativa de HIT es **menor** que en climas cálidos: el aire "
        "de altura de Atacama enfría la celda (≈58–62 °C), reduciendo el castigo "
        "térmico que favorece a HIT.\n\n"
        "• Aun así HIT gana en PR, energía y facturación en ambos escenarios de "
        "recurso (emulado y real).",
        accent=ACCENT_GREEN, body_size=16.5)


def s_cumplimiento(slide, n):
    add_title(slide, "Síntesis: Qué Pidió la Tarea 2 y Cómo se Resolvió")
    add_panel(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(5.7),
              "Mapa de cumplimiento — requisito de la pauta → evidencia en esta presentación",
              accent=ACCENT_GREEN)
    headers = ["Requisito (pauta Tarea 2)", "Cómo se resolvió", "Evidencia", "Resultado clave"]
    rows = [
        ["Recurso solar: G[_poa_] + óptica",
         "Perez + IAM físico + factor espectral",
         "Lám. %d–%d · Anexos I–II" % (REF['f2_flujo'], REF['optica']),
         "POA efectiva 5-min, 2026"],
        ["Perfil térmico T[_c_] — ec. (1)",
         "Sandia/SAPM, viento 1 m/s",
         "Lám. %d · Anexo III" % REF['dia_tipico'],
         "T[_c_] pico 70–73 °C"],
        ["SDM 5 parámetros De Soto — ec. (2)",
         "`minimize` + bounds físicos en SRC",
         "Lám. %d · Anexos V–VI" % REF['f3_result'],
         "R² ≈ 0.99 vs medición"],
        ["Discusión R[_s_] – n (puntos conflictuales)",
         "Mal condicionamiento + mitigación",
         "Lám. %d" % REF['rs_n'],
         "n[_I_] = 1.20; bounds activos"],
        ["PR anual 2026 — ec. (3)",
         "Σ P[_mp_] / Σ P[_STC_]·G/1000 (IEC 61724-1)",
         "Lám. %d, %d · Anexo VII" % (REF['veredicto_eco'], REF['recurso_real']),
         "HIT 84.18 % · m-Si 81.61 %"],
        ["Supuestos propuestos y justificados",
         "tilt 22.9° = lat · azimut 0° N · viento 1 m/s · albedo 0.20",
         "Lám. %d y %d" % (REF['emu_impl'], REF['limitaciones']),
         "Declarados + limitaciones"],
    ]
    add_table(slide, MARGIN + Inches(0.25), CONTENT_Y + Inches(0.60),
              SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(4.85),
              headers, rows, col_ratios=[2.3, 2.4, 1.5, 1.7],
              fs_head=15.0, fs_body=13.5)


def s_limitaciones(slide, n):
    add_title(slide, "d) Limitaciones del Estudio y Alcance de Validez")
    H = Inches(5.45)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Resuelto en esta versión",
        "• **Modificadores ópticos APLICADOS:** IAM físico + factor espectral en "
        "Fase 2 " + A('anexo2') + " (≈3 %% de pérdida óptica ya penalizada).\n\n"
        "• **Recurso REAL de Atacama:** track con PVGIS TMY (Lám. %d) entrega "
        "energía y economía absolutas (POA ≈ 2,810 kWh/m²·año).\n\n"
        "• **Doble cobertura corregida:** ventana de 12 meses contiguos → cada mes "
        "se cubre una sola vez (antes jul–sep duplicaban temporada).\n\n"
        "• **Albedo declarado 0.20** y energía integrada al paso real de 5 min."
        % REF['recurso_real'],
        accent=ACCENT_GREEN, body_size=15.5, line_space=1.05)
    panel_with_body(slide, MARGIN + COL_W + GAP, CONTENT_Y, COL_W, H,
        "Limitaciones que permanecen",
        "• **Validación eléctrica con datos de Florida:** el R² ≈ 0.99 se valida "
        "contra mediciones de Cocoa; Atacama no tiene un módulo PV de referencia "
        "medido, por lo que el track real usa TMY satelital (PVGIS), no in-situ.\n\n"
        "• **α[_Isc_] de literatura:** la regresión experimental falló y se usó "
        "+0.05 %/°C — impacto de segundo orden (la corriente la domina G).\n\n"
        "• **P[_STC_] del SDM** (50.1 / 238.0 W) difiere ~7–9 % del promedio "
        "empírico (46.7 / 218.4 W); criterio consistente entre tecnologías.\n\n"
        "• **Sin soiling, BOS ni mismatch:** PR de módulo, no de planta; modelo "
        "de un diodo.",
        accent=ACCENT_BLUE, body_size=15.5, line_space=1.05)


def s_conclusiones(slide, n):
    add_title(slide, "e) Conclusiones y propuestas de trabajos futuros")
    H = Inches(5.45)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H,
        "Conclusiones del Estudio",
        "1. **Fidelidad metodológica:** la emulación por traslación temporal "
        "(+6 meses, año 2026) mantuvo la coherencia de solsticios y permitió "
        "simular Atacama con datos experimentales reales.\n\n"
        "2. **Pipeline de datos robusto:** la ingesta en streaming resolvió "
        "archivos irregulares de ~110 MB conservando ~96 % de la ventana de 12 meses.\n\n"
        "3. **Ajuste robusto:** el SDM de 5 parámetros, con IAM y corrección "
        "espectral aplicados, reprodujo la potencia medida con R² ≈ 0.99.\n\n"
        "4. **HIT vencedor:** +2.57 pts de PR validado (84.18 % vs 81.61 %); con "
        "el recurso real de Atacama, +3,294 MWh y +USD 148k/año en 100 MWp.",
        accent=ACCENT_GOLD, body_size=16.0, line_space=1.05)
    panel_with_body(slide, MARGIN + COL_W + GAP, CONTENT_Y, COL_W, H,
        "Propuestas de Trabajos Futuros",
        "• **Validación in-situ:** contrastar contra un módulo PV medido en "
        "Atacama (estación de referencia), no solo TMY satelital.\n\n"
        "• **Soiling desértico:** acumulación de polvo, factor crítico de "
        "pérdida en Atacama.\n\n"
        "• **Pérdidas de planta (BOS):** cableado, inversor y mismatch para pasar "
        "de PR de módulo a PR de planta.\n\n"
        "• **Bifacialidad:** aporte del albedo desértico en la cara posterior.\n\n"
        "• **Modelo de doble diodo:** capturar recombinación no ideal a baja "
        "irradiancia.",
        accent=ACCENT_BLUE, body_size=16.0, line_space=1.05)


def s_referencias(slide, n):
    add_title(slide, "Referencias")
    panel_with_body(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(5.5),
        "Referencias Bibliográficas Estructuradas",
        '1. De Soto, W., Klein, S.A., Beckman, W.A. (2006). "Improvement and '
        'validation of a model for photovoltaic array performance." Solar Energy '
        '80, 78–88.\n\n'
        '2. King, D.L., Boyson, W.E., Kratochvil, J.A. (2004). "Photovoltaic '
        'Array Performance Model." Sandia Report SAND2004-3535.\n\n'
        '3. Marion, B. et al. (2014). "New Data Set for Validating PV Module '
        'Performance Models." NREL/CP-5J00-61610 — origen del dataset Cocoa.\n\n'
        '4. Perez, R. et al. (1990). "Modeling daylight availability and '
        'irradiance components from direct and global irradiance." Solar Energy '
        '44(5), 271–289.\n\n'
        '5. Holmgren, W.F., Hansen, C.W., Mikofski, M.A. (2018). "pvlib python: '
        'a python package for modeling solar energy systems." JOSS 3(29), 884.\n\n'
        '6. IEC 61724-1 — Photovoltaic system performance, Part 1: Monitoring.\n\n'
        'Agradecimientos al Departamento de Electrotecnia de la UTFSM.',
        accent=ACCENT_ORANGE, body_size=15.5, line_space=1.08)


def s_anexo_tools(slide, n):
    add_title(slide, "Anexo · Librerías y Herramientas de Código Utilizadas")
    H1 = Inches(2.62)
    gap_v = Inches(0.22)
    X2 = MARGIN + COL_W + GAP
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H1,
        "Lenguaje y núcleo científico",
        "• **Python 3.12** en entorno virtual (`venv`).\n"
        "• **pvlib 0.15** — modelado fotovoltaico (transposición, IAM, espectral, "
        "térmico, modelo de diodo, lectura PVGIS).\n"
        "• **NumPy 2.4** — cálculo numérico vectorizado.\n"
        "• **pandas 3.0** — series temporales, `resample`, limpieza.\n"
        "• **SciPy** — `optimize.minimize`/`fsolve` (ajuste de 5 parámetros).",
        accent=ACCENT_BLUE, body_size=14.0, line_space=1.06)
    panel_with_body(slide, MARGIN, CONTENT_Y + H1 + gap_v, COL_W, H1,
        "Funciones pvlib utilizadas (por fase)",
        "• `irradiance.get_total_irradiance` (Perez), `aoi`, `get_extra_radiation`.\n"
        "• `iam.physical` + `iam.marion_diffuse` (IAM directo y difuso).\n"
        "• `spectrum.spectral_factor_firstsolar`, `atmosphere.gueymard94_pw` + airmass.\n"
        "• `temperature.sapm_cell` (modelo térmico Sandia).\n"
        "• `pvsystem.calcparams_desoto` · `singlediode` · `i_from_v` (Lambert W).\n"
        "• `iotools.get_pvgis_tmy` · `location.Location`.",
        accent=ACCENT_GOLD, body_size=14.0, line_space=1.06)
    panel_with_body(slide, X2, CONTENT_Y, COL_W, H1,
        "Visualización y generación de entregables",
        "• **Matplotlib** — gráficos (modo oscuro, `mathtext` para ecuaciones).\n"
        "• **python-pptx** — generación automática del archivo PPTX.\n"
        "• **Pillow (PIL)** — recorte e inserción de imágenes y ecuaciones.\n"
        "• **pywin32** (`win32com`) — exporta cada lámina a PNG vía PowerPoint.",
        accent=ACCENT_GREEN, body_size=14.0, line_space=1.06)
    panel_with_body(slide, X2, CONTENT_Y + H1 + gap_v, COL_W, H1,
        "Utilidades y datos externos",
        "• **python-dateutil** (`relativedelta`) — desfase estacional de +6 meses.\n"
        "• **stdlib:** `csv`, `json`, `os`, `re`, `datetime`, `hashlib`, `math`.\n"
        "• **Dataset NREL Cocoa** — mediciones experimentales (11 módulos).\n"
        "• **PVGIS TMY** (JRC, Comisión Europea, base SARAH) — recurso real de "
        "Atacama (Fase 8).",
        accent=ACCENT_ORANGE, body_size=14.0, line_space=1.06)


def s_anexo1(slide, n):
    add_title(slide, "Anexo I: Transposición de Irradiancia al Plano del Arreglo (Perez)")
    H = Inches(5.45)
    add_panel(slide, MARGIN, CONTENT_Y, COL_W, H,
              "Forma general de transposición POA", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.62), COL_W - Inches(0.44), Inches(0.6),
             "Irradiancia incidente total en el plano inclinado del panel:", size=17.0)
    add_eq(slide,
           r"G_{poa}=G_b\,R_{beam}+G_{d,perez}+G\,\rho\,\left(\frac{1-\cos(\beta)}{2}\right)",
           MARGIN + Inches(0.35), CONTENT_Y + Inches(1.25), Inches(0.62), max_width=Inches(5.7))
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(2.3), COL_W - Inches(0.44), Inches(2.9),
             "Donde:\n"
             "• G[_b_] / G[_d_] / G: irradiancia directa, difusa y global (W/m²).\n"
             "• β: inclinación del panel = 22.91° (latitud del sitio).\n"
             "• ρ: albedo = 0.20 (suelo desértico, declarado explícitamente).\n"
             "• R[_beam_]: factor geométrico de transposición directa.\n"
             "• G[_d,perez_]: difusa anisotrópica de Perez (circumsolar +\n"
             "  brillo de horizonte), evaluada por `pvlib.get_total_irradiance`.",
             size=16.0, color=ACCENT_GOLD, line_space=1.10)
    X2 = MARGIN + COL_W + GAP
    panel_with_body(slide, X2, CONTENT_Y, COL_W, H,
        "Implementación en el código (Fase 2)",
        "• `location.get_solarposition(times)` → cenit y azimut aparentes del "
        "sol registro a registro.\n\n"
        "• `irradiance.get_extra_radiation(times)` → irradiancia extraterrestre "
        "para el modelo anisotrópico.\n\n"
        "• `irradiance.get_total_irradiance(tilt, az, dni, ghi, dhi, ..., "
        "albedo=0.20, model='perez')` → directa + difusa Perez + reflejada.\n\n"
        "• Salida: `poa_global` (banda ancha) → entra a los modificadores "
        "ópticos IAM y espectral del Anexo II para dar la irradiancia efectiva.",
        accent=ACCENT_GOLD, body_size=16.0)


def s_anexo2(slide, n):
    add_title(slide, "Anexo II: Modificadores Ópticos APLICADOS (IAM y Espectral)")
    H = Inches(4.4)
    add_panel(slide, MARGIN, CONTENT_Y, COL_W, H,
              "Modificador por Ángulo de Incidencia (IAM)", accent=ACCENT_BLUE)
    add_eq(slide, r"K_{\tau\alpha}(\theta)=\frac{\tau(\theta)}{\tau(0)}",
           MARGIN + Inches(0.35), CONTENT_Y + Inches(0.75), Inches(0.55), max_width=Inches(5.5))
    add_eq(slide,
           r"\tau(\theta)=e^{-\frac{KL}{\cos(\theta_r)}}\left[1-\frac{1}{2}\left(\frac{\sin^2(\theta_r-\theta)}{\sin^2(\theta_r+\theta)}+\frac{\tan^2(\theta_r-\theta)}{\tan^2(\theta_r+\theta)}\right)\right]",
           MARGIN + Inches(0.35), CONTENT_Y + Inches(1.55), Inches(0.75), max_width=Inches(5.6))
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(2.7), COL_W - Inches(0.44), Inches(1.5),
             "Ley de Snell y Bouguer [King et al., 2004]: θ[_r_] = arcsin(sin θ / n), con "
             "n = 1.526 (vidrio), K = 4 m⁻¹, L = 2 mm.",
             size=16.5, color=ACCENT_GOLD, line_space=1.15)
    X2 = MARGIN + COL_W + GAP
    add_panel(slide, X2, CONTENT_Y, COL_W, H,
              "Modificador Espectral por Masa de Aire (AM)", accent=ACCENT_GOLD)
    add_eq(slide,
           r"\frac{M}{M_{ref}}=a_0+a_1 AM+a_2 AM^2+a_3 AM^3+a_4 AM^4",
           X2 + Inches(0.35), CONTENT_Y + Inches(0.8), Inches(0.55), max_width=Inches(5.6))
    add_eq(slide,
           r"AM=\frac{1}{\cos(\theta_z)+0.5057\,(96.08-\theta_z)^{-1.634}}",
           X2 + Inches(0.35), CONTENT_Y + Inches(1.6), Inches(0.65), max_width=Inches(5.6))
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(2.7), COL_W - Inches(0.44), Inches(1.5),
             "Implementado con `pvlib.spectrum.spectral_factor_firstsolar` (agua precipitable de "
             "humedad+T y masa de aire absoluta), módulo c-Si ('monosi').",
             size=16.5, color=ACCENT_GOLD, line_space=1.15)
    add_caption(slide,
        "**Estado en este estudio: EJECUTADO en la Fase 2.** El IAM directo usa `pvlib.iam.physical` "
        "(n=1.526, K=4, L=2 mm) y la difusa los factores de Marion; el factor espectral usa First Solar. "
        "La irradiancia efectiva = POA·IAM·M alimenta el SDM. Efecto en Atacama: ≈3 % de pérdida óptica y "
        "espectral casi neutra (cielo seco y limpio).",
        top_in=5.80, h_in=1.30, accent=ACCENT_GREEN)


def s_anexo3(slide, n):
    add_title(slide, "Anexo III: Ecuaciones del Perfil Térmico (Modelo Sandia SAPM)")
    add_panel(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(1.95),
              "Ecuación de Sandia para Temperatura de Celda", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.58),
             SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.4),
             "La temperatura de la celda se modela a partir del equilibrio térmico dinámico:", size=18.0)
    add_eq(slide,
           r"T_c = G_{poa}\cdot e^{a+b\,v_w} + T_a + \left(\frac{G_{poa}}{1000}\right)\Delta T",
           Inches(4.0), CONTENT_Y + Inches(1.05), Inches(0.7), max_width=Inches(6.0))
    Y2 = CONTENT_Y + Inches(2.15)
    H2 = Inches(3.3)
    panel_with_body(slide, MARGIN, Y2, COL_W, H2,
         "Glosario de Variables",
         "• T[_c_] / T[_a_]: temperatura de celda y ambiente (°C).\n"
         "• G[_poa_]: irradiancia total en el plano del panel (W/m²).\n"
         "• v[_w_]: velocidad de viento (fijada en 1 m/s).\n"
         "• a, b, ΔT: parámetros empíricos del encapsulado/montaje.",
         accent=ACCENT_BLUE, body_size=18.0)
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
         "Coeficientes Empíricos Utilizados",
         "• m-Si — `open_rack_glass_polymer`: a = −3.56 · b = −0.075 · ΔT = 3.0 °C\n\n"
         "• HIT — `open_rack_glass_glass`: a = −3.47 · b = −0.059 · ΔT = 3.0 °C\n\n"
         "**Fuente:** `pvlib.temperature.TEMPERATURE_MODEL_PARAMETERS['sapm']` — "
         "base interna validada por Sandia NL [King et al., 2004].",
         accent=ACCENT_GOLD, body_size=15.0)


def s_anexo4(slide, n):
    add_title(slide, "Anexo IV: Ecuaciones del Modelo Eléctrico de un Diodo (SDM)")
    add_panel(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(1.95),
               "Ecuación Trascendental del Circuito Equivalente", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.58),
              SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.6),
              "Corriente implícita del modelo de diodo simple con resistencias parásitas "
              "(resuelta con la función W de Lambert en `pvlib.singlediode`):", size=16.0)
    add_eq(slide,
            r"I = I_L - I_0\left[\exp\left(\frac{V+I R_s}{a}\right)-1\right]-\frac{V+I R_s}{R_{sh}}",
            Inches(3.8), CONTENT_Y + Inches(1.1), Inches(0.62), max_width=Inches(6.2))
    Y2 = CONTENT_Y + Inches(2.15)
    H2 = Inches(3.3)
    add_panel(slide, MARGIN, Y2, COL_W, H2,
               "Factor de Idealidad Térmico (a)", accent=ACCENT_GOLD)
    add_eq(slide, r"a=\frac{N_s\,n_I\,k\,T_c}{q}",
            MARGIN + Inches(0.45), Y2 + Inches(0.75), Inches(0.62), max_width=Inches(3.2))
    add_text(slide, MARGIN + Inches(0.22), Y2 + Inches(1.7), COL_W - Inches(0.44), Inches(1.4),
              "• N[_s_]: celdas en serie (m-Si: 36, HIT: 72).\n"
              "• n[_I_]: factor de idealidad del diodo, acotado a [0.6, 2.4] "
              "(ajustado: 1.20 en ambas tecnologías).",
              size=18.0, color=ACCENT_GOLD, line_space=1.15)
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
        "Glosario de Parámetros Físicos",
        "• I / V: corriente y voltaje de salida (A, V).\n"
        "• I[_L_]: corriente fotogenerada (A).\n"
        "• I[_0_]: corriente de saturación inversa de la unión P-N (A).\n"
        "• R[_s_] / R[_sh_]: resistencias parásitas serie y shunt (Ω).\n"
        "• k / q / T[_c_]: constante de Boltzmann, carga elemental y "
        "temperatura de celda (K).",
        accent=ACCENT_BLUE, body_size=18.0)


def s_anexo5(slide, n):
    add_title(slide, "Anexo V: Extracción en STC — Ajuste Numérico Real")
    H = Inches(5.45)
    add_panel(slide, MARGIN, CONTENT_Y, COL_W, H,
              "Normalización de mediciones a SRC", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.6), COL_W - Inches(0.44), Inches(0.7),
             "Cada punto en la ventana 900–1100 W/m² se traslada a 1000 W/m² y 25 °C:",
             size=18.0)
    add_eq(slide, r"I_{sc}^{SRC}=I_{sc}\,\frac{1000}{G}+\alpha_{Isc}\,(25-T_c)",
           MARGIN + Inches(0.4), CONTENT_Y + Inches(1.3), Inches(0.55), max_width=Inches(5.4))
    add_eq(slide, r"V_{oc}^{SRC}=V_{oc}+\beta_{Voc}\,(25-T_c)",
           MARGIN + Inches(0.4), CONTENT_Y + Inches(2.05), Inches(0.55), max_width=Inches(5.4))
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(2.85), COL_W - Inches(0.44), Inches(2.3),
             "• β[_Voc_] (medido): regresión lineal V[_oc_] vs T[_c_] (ventana 800–1200 W/m²).\n"
             "• α[_Isc_] (literatura): la regresión resultó negativa → resguardo físico "
             "+0.05 %/°C de literatura.\n"
             "• Promedios de la ventana → I[_sc,ref_], V[_oc,ref_], I[_mp,ref_], "
             "V[_mp,ref_].",
             size=17.5, color=ACCENT_GOLD, line_space=1.2)
    X2 = MARGIN + COL_W + GAP
    add_panel(slide, X2, CONTENT_Y, COL_W, H,
              "Ajuste de los 5 parámetros (lo que ejecuta el código)", accent=ACCENT_GOLD)
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(0.6), COL_W - Inches(0.44), Inches(0.65),
             "`scipy.optimize.minimize` sobre el residuo cuadrático de las 3 condiciones "
             "del SDM en SRC:", size=17.5)
    add_eq(slide, r"\min_{I_L,I_0,a,R_s,R_{sh}}\;e_{V_{oc}}^2+e_{I_{sc}}^2+e_{MPP}^2",
           X2 + Inches(0.45), CONTENT_Y + Inches(1.35), Inches(0.6), max_width=Inches(5.2))
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(2.25), COL_W - Inches(0.44), Inches(2.9),
             "• e[_Voc_]: corriente nula en circuito abierto.\n"
             "• e[_Isc_]: corriente de cortocircuito reproducida.\n"
             "• e[_MPP_]: corriente correcta en (V[_mp_], I[_mp_]).\n"
             "**Bounds físicos:** R[_s_] ∈ [0.001, 2] Ω · R[_sh_] ∈ [100, 10⁴] Ω "
             "· a ∈ [0.5, 2]·a₀ · I[_L_] ± 10 % · I[_0_] ∈ [10⁻¹², 10⁻⁵] A.\n"
             "**Inicialización analítica:** a₀ = N[_s_]·1.2·kT/q; "
             "I[_0,0_] = I[_sc_]·exp(−V[_oc_]/a₀).\n"
             "3 ecuaciones, 5 incógnitas: la unicidad la dan bounds + "
             "inicialización (R[_s_], R[_sh_] quedan ≈ iniciales).",
             size=15.0, color=ACCENT_GOLD, line_space=1.12)


def s_anexo6(slide, n):
    add_title(slide, "Anexo VI: Escalamiento a Condiciones de Operación (De Soto)")
    H = Inches(5.45)
    add_panel(slide, MARGIN, CONTENT_Y, COL_W, H,
              "Corrientes y Factor de Idealidad", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.72), COL_W - Inches(0.44), Inches(0.35),
             "1. Corriente fotogenerada (I[_L_]):", size=17.5)
    add_eq(slide, r"I_L=\frac{S}{S_{ref}}\left[I_{L,ref}+\alpha_{Isc}(T_c-T_{ref})\right]",
           MARGIN + Inches(0.4), CONTENT_Y + Inches(1.10), Inches(0.55), max_width=Inches(5.3))
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(1.85), COL_W - Inches(0.44), Inches(0.35),
             "2. Corriente de saturación inversa (I[_0_]):", size=17.5)
    add_eq(slide,
           r"\frac{I_0}{I_{0,ref}}=\left(\frac{T_c}{T_{ref}}\right)^3\exp\left[\frac{E_{g,ref}}{k\,T_{ref}}-\frac{E_g}{k\,T_c}\right]",
           MARGIN + Inches(0.4), CONTENT_Y + Inches(2.28), Inches(0.65), max_width=Inches(5.3))
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(3.15), COL_W - Inches(0.44), Inches(0.35),
             "3. Factor de idealidad térmico (a):", size=17.5)
    add_eq(slide, r"\frac{a}{a_{ref}}=\frac{T_c}{T_{ref}}",
           MARGIN + Inches(0.4), CONTENT_Y + Inches(3.60), Inches(0.55), max_width=Inches(3.0))
    X2 = MARGIN + COL_W + GAP
    add_panel(slide, X2, CONTENT_Y, COL_W, H,
              "Resistencias y Energía de Bandgap", accent=ACCENT_GOLD)
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(0.72), COL_W - Inches(0.44), Inches(0.35),
             "4. Resistencia shunt (R[_sh_]):", size=17.5)
    add_eq(slide, r"R_{sh}=R_{sh,ref}\left(\frac{S_{ref}}{S}\right)",
           X2 + Inches(0.4), CONTENT_Y + Inches(1.10), Inches(0.6), max_width=Inches(4.0))
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(1.90), COL_W - Inches(0.44), Inches(0.35),
             "5. Energía de bandgap del Silicio (E[_g_]):", size=17.5)
    add_eq(slide, r"\frac{E_g}{E_{g,ref}}=1-0.0002677\,(T_c-T_{ref})",
           X2 + Inches(0.4), CONTENT_Y + Inches(2.30), Inches(0.55), max_width=Inches(4.8))
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(3.15), COL_W - Inches(0.44), Inches(2.2),
             "Suposiciones físicas justificadas:\n"
             "• R[_s_] constante: R[_s_] = R[_s,ref_] (validación NIST < 2 %).\n"
             "• R[_sh_] ∝ 1/S: empírica, dominante a baja irradiancia.\n"
             "• E[_g,ref_] = 1.121 eV (Si, 25 °C); S = G[_poa_] (Anexo II).\n"
             "Implementado por `pvlib.pvsystem.calcparams_desoto`.",
             size=15.5, color=ACCENT_GOLD, line_space=1.12)


def s_anexo7(slide, n):
    add_title(slide, "Anexo VII: Ecuaciones para el Cálculo del Performance Ratio")
    add_panel(slide, MARGIN, CONTENT_Y, SLIDE_W - 2 * MARGIN, Inches(1.95),
              "Ecuación de Rendimiento Global del PR (IEC 61724-1)", accent=ACCENT_BLUE)
    add_text(slide, MARGIN + Inches(0.22), CONTENT_Y + Inches(0.72),
             SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.4),
             "Eficiencia neta del sistema frente a su comportamiento ideal en condiciones estándar:",
             size=18.0)
    add_eq(slide,
           r"PR=\sum_{t=1}^{N} P_{mp,t}(G_t,T_{c,t})\;/\;\sum_{t=1}^{N} P_{STC}\cdot\frac{G_{poa,t}}{G_{ref}}",
           Inches(3.9), CONTENT_Y + Inches(1.05), Inches(0.78), max_width=Inches(6.0))
    Y2 = CONTENT_Y + Inches(2.15)
    H2 = Inches(3.3)
    panel_with_body(slide, MARGIN, Y2, COL_W, H2,
        "Glosario de Variables",
        "• P[_mp,t_]: potencia MPP del SDM en el registro t (cadencia 5 min).\n"
        "• P[_STC_]: potencia nominal del **mismo SDM** evaluado a 1000 W/m² y "
        "25 °C — m-Si: 50.1 W · HIT: 238.0 W (criterio consistente).\n"
        "• G[_poa,t_]: irradiancia instantánea en el plano del panel (W/m²).\n"
        "• G[_ref_]: irradiancia de referencia STC (1000 W/m²).",
        accent=ACCENT_GOLD, body_size=15.0)
    panel_with_body(slide, MARGIN + COL_W + GAP, Y2, COL_W, H2,
        "Pérdidas penalizadas por el PR",
        "• **Térmicas:** caídas por elevada T[_c_] — dominantes en Atacama.\n"
        "• **Óhmicas:** pérdidas resistivas en R[_s_].\n"
        "• **De bajo G:** comportamiento no lineal a baja irradiancia (R[_sh_]).\n"
        "• NO penaliza pérdidas ópticas/espectrales (no modeladas) ni de "
        "planta (cableado, inversor, soiling).",
        accent=ACCENT_BLUE, body_size=18.0)


def s_anexo8(slide, n):
    add_title(slide, "Anexo VIII: Diccionario de Columnas del Dataset NREL Cocoa")
    H1 = Inches(2.58)
    panel_with_body(slide, MARGIN, CONTENT_Y, COL_W, H1,
        "Identificación y tiempo",
        "• `Time Stamp` (hora local estándar, paso 5 min) — col. 0.\n"
        "• Metadatos (líneas 1–2): módulo, ciudad, estado, zona horaria, "
        "latitud, longitud, altitud, tilt, azimut.",
        accent=ACCENT_BLUE, body_size=15.5)
    panel_with_body(slide, MARGIN, CONTENT_Y + H1 + Inches(0.25), COL_W, H1,
        "Eléctricas (usadas → índices fijos)",
        "• `Isc` (5), `Pmp` (7), `Imp` (9), `Vmp` (11), `Voc` (13) — cada una "
        "con su columna de incertidumbre (%).\n"
        "• `FF` (15) y cola de n pares (I, V) crudos — descartadas en la "
        "ingesta.",
        accent=ACCENT_GREEN, body_size=15.5)
    X2 = MARGIN + COL_W + GAP
    panel_with_body(slide, X2, CONTENT_Y, COL_W, H1,
        "Meteorológicas (usadas → índices fijos)",
        "• `Dry bulb temperature` (20), `Atmospheric pressure` (24).\n"
        "• `DNI` (27), `GHI` (30), `DHI` (33) — c/u con incertidumbre y desv. "
        "estándar de muestras de 1 s.\n"
        "• `POA CMP22` (1): referencia para verificar el recurso.",
        accent=ACCENT_GREEN, body_size=15.5)
    panel_with_body(slide, X2, CONTENT_Y + H1 + Inches(0.25), COL_W, H1,
        "Calidad y operación (no usadas)",
        "• `Solar QA residual`: cierre Direct·cos(z) + Difusa − Global.\n"
        "• `PV module soiling derate`, lluvia, humedad relativa, T dorso, "
        "T gabinete MT5, horarios de mantención (`99:99` = sin mantención) y "
        "n.º de pares I-V.",
        accent=ACCENT_ORANGE, body_size=15.5)


def s_optica(slide, n):
    add_title(slide, "Procedimiento: Modificadores Ópticos APLICADOS — IAM y Espectral (Fase 2)")
    full_width_image(slide, IMG['optica'], max_w=12.2, top_in=1.28, max_bottom=5.15)
    add_caption(slide,
        "**Pérdidas ópticas y espectrales ahora ejecutadas en el pipeline:** sobre la POA de Perez (albedo "
        "**0.20**) se aplica el **IAM físico** de Snell-Bouguer (`pvlib.iam.physical`, n=1.526, K=4 m⁻¹, L=2 mm) "
        "a la directa y los factores de Marion a la difusa, y el **factor espectral** First Solar "
        "(`spectral_factor_firstsolar`, agua precipitable + masa de aire). El resultado es la **irradiancia "
        "efectiva** que fotogenera I[_L_] en la Fase 4. Efecto en Atacama: pérdida óptica ≈ **3.0 %%** "
        "(ángulos rasantes) y espectral ≈ **neutra** (cielo seco y limpio) — ecuaciones en %s." % A('anexo2'),
        top_in=5.25, h_in=1.75, accent=ACCENT_GOLD, size=15.5)


def s_recurso_real(slide, n):
    add_title(slide, "Resultado: Recurso REAL de Atacama (PVGIS TMY) — Magnitudes Absolutas")
    TW = (COL_W - Inches(0.4)) / 2
    stat_tile(slide, MARGIN, CONTENT_Y, TW, Inches(1.5),
              "2,596", "GHI real anual del sitio\n(kWh/m²·año, PVGIS TMY)", accent=ACCENT_BLUE)
    stat_tile(slide, MARGIN + TW + Inches(0.4), CONTENT_Y, TW, Inches(1.5),
              "2,810", "POA en plano inclinado\n(kWh/m²·año, albedo 0.20)", accent=ACCENT_GOLD)
    stat_tile(slide, MARGIN, CONTENT_Y + Inches(1.70), TW, Inches(1.5),
              "2,389", "Yield HIT\n(kWh/kWp·año)", accent=ACCENT_GOLD)
    stat_tile(slide, MARGIN + TW + Inches(0.4), CONTENT_Y + Inches(1.70), TW, Inches(1.5),
              "2,356", "Yield m-Si\n(kWh/kWp·año)", accent=ACCENT_BLUE)
    panel_with_body(slide, MARGIN, CONTENT_Y + Inches(3.35), COL_W, Inches(2.25),
        "Lectura con recurso real (planta 100 MWp)",
        "• **Fuente:** PVGIS (JRC, UE), base satelital SARAH — TMY horario de "
        "San Pedro de Atacama.\n"
        "• HIT gana **+1.17 pts de PR**; el aire frío de altura modera la celda "
        "a ≈ 58–62 °C.\n"
        "• Δyield 33 kWh/kWp → **+3,294 MWh, +USD 148k/año** (100 MWp) %s." % A('anexo13'),
        accent=ACCENT_GREEN, body_size=13.5, line_space=1.04)
    X2 = MARGIN + COL_W + GAP
    add_panel(slide, X2, CONTENT_Y, COL_W, Inches(5.45),
              "PR mensual con recurso real (PVGIS TMY)", accent=ACCENT_BLUE)
    add_image(slide, IMG['pr_real'], X2 + Inches(0.22), CONTENT_Y + Inches(1.35),
              width=Inches(5.85), height=Inches(2.78))
    add_text(slide, X2 + Inches(0.22), CONTENT_Y + Inches(4.35), COL_W - Inches(0.44), Inches(0.95),
             "El recurso real ≈ 2.3× el emulado (Florida): la energía absoluta se "
             "dispara, y HIT mantiene la ventaja aunque el sitio sea térmicamente más benigno.",
             size=13.5, color=TEXT_GREY, line_space=1.05)


BUILDERS = {
    'portada': s_portada, 'intro': s_intro, 'justificacion': s_justificacion,
    'marco': s_marco, 'pipeline': s_pipeline, 'datos': s_datos,
    'ingesta': s_ingesta, 'embudo': s_embudo,
    'emu_concepto': s_emu_concepto, 'emu_impl': s_emu_impl,
    'poa_medida': s_poa_medida, 'f2_flujo': s_f2_flujo,
    'dia_tipico': s_dia_tipico, 'f3_flujo': s_f3_flujo,
    'f3_result': s_f3_result, 'f4_flujo': s_f4_flujo,
    'optica': s_optica, 'recurso_real': s_recurso_real,
    'validacion': s_validacion, 'rs_n': s_rs_n,
    'perdidas_pr': s_perdidas_pr, 'veredicto_eco': s_veredicto_eco,
    'cumplimiento': s_cumplimiento, 'limitaciones': s_limitaciones,
    'conclusiones': s_conclusiones, 'referencias': s_referencias,
    'anexo_tools': s_anexo_tools,
    'anexo1': s_anexo1, 'anexo2': s_anexo2,
    'anexo3': s_anexo3, 'anexo4': s_anexo4, 'anexo5': s_anexo5,
    'anexo6': s_anexo6, 'anexo7': s_anexo7, 'anexo8': s_anexo8,
    'anexo9': s_anexo9, 'anexo10': s_anexo10, 'anexo11': s_anexo11,
    'anexo12': s_anexo12, 'anexo13': s_anexo13,
}

# ============================================================
# NOTAS DEL EXPOSITOR — (segundos asignados, mensaje clave)
# Presupuesto total láminas 1-24: ~19.5 min (colchón ~0.5 min en 20 min).
# ============================================================
NOTES = {
    'portada': (20, "Saludo. Somos consultores técnicos: ¿qué tecnología PV conviene "
                "instalar en Atacama en 2026? Respuesta basada en datos medidos, no en catálogo."),
    'intro': (60, "Atacama = mejor recurso solar del mundo, pero las celdas operan a 65-73 °C "
              "y el calor degrada la potencia. Objetivo: cuantificar qué tecnología resiste mejor."),
    'justificacion': (50, "De las 5 familias del dataset elegimos los 2 extremos térmicos "
                      "comercialmente viables: m-Si (estándar, sensible) vs HIT (premium, tolerante)."),
    'marco': (50, "Modelos publicados encadenados: Perez (POA), Sandia (Tc), De Soto SDM (eléctrico) y su "
              "escalamiento, MÁS los modificadores ópticos IAM y espectral, ahora aplicados en Fase 2."),
    'pipeline': (60, "Mapa del trabajo: 5 fases automatizadas desde 1.2 GB de CSV crudos hasta el PR. "
                 "Esta lámina es la guía de las secciones siguientes."),
    'datos': (50, "La base: 11 módulos medidos en simultáneo cada 5 min por NREL. Usamos 11 columnas "
              "de 43; cada registro es una curva I-V real con su meteorología."),
    'ingesta': (45, "Los archivos no se pueden abrir con read_csv: fila variable + sentinelas -9999. "
                "Lector streaming propio → ~96 % de la ventana de 12 meses."),
    'embudo': (45, "Trazabilidad completa: de 420 mil curvas brutas a ~64 mil registros simulados. "
               "Cada filtro tiene número y razón."),
    'emu_concepto': (60, "Idea clave de la emulación: +6 meses alinea los solsticios entre hemisferios. "
                     "Sin esto, simularíamos un Atacama con verano en julio."),
    'emu_impl': (45, "Implementación: reescribir metadatos de sitio + desfase de fechas. Las magnitudes "
                 "físicas NO se tocan — límite declarado del método."),
    'poa_medida': (40, "Verificación previa: las 11 curvas POA del piranómetro patrón se superponen → "
                   "los sensores son consistentes y la comparación entre módulos es justa."),
    'f2_flujo': (45, "Fase 2: Perez transpone GHI/DNI/DHI al plano del panel; Sandia estima la "
                 "temperatura de celda con viento conservador de 1 m/s."),
    'optica': (45, "Ahora SÍ aplicamos los modificadores ópticos: IAM físico (pérdida ~3 % en ángulos "
               "rasantes) y factor espectral (neutro en el cielo seco de Atacama). La irradiancia efectiva "
               "alimenta el modelo eléctrico — antes era trabajo futuro, ahora está ejecutado."),
    'dia_tipico': (40, "Para no expertos: en un día despejado la celda sigue al sol y supera los 60 °C. "
                   "Cada grado sobre 25 °C cuesta potencia — eso es lo que vamos a medir."),
    'f3_flujo': (45, "Fase 3: coeficientes térmicos por regresión (β medido; α de literatura porque "
                 "la regresión salió negativa) y ajuste de los 5 parámetros con bounds físicos."),
    'f3_result': (60, "Parámetros extraídos: el I0 de HIT es 10 veces menor (menos recombinación). "
                  "Honestidad: Rs y Rsh quedan en su inicialización — lo discutimos en la lámina Rs-n."),
    'f4_flujo': (40, "Fase 4: para cada registro 5-min se escala el circuito con De Soto y se resuelve "
                 "la potencia máxima con Lambert W. De ahí sale el PR anual."),
    'validacion': (55, "Antes de comparar tecnologías, validamos el modelo: R² 0.99 contra la potencia "
                   "medida por el trazador. Mismo G y Tc en ambos ejes — valida el circuito, no compara sitios."),
    'rs_n': (60, "Punto conflictual exigido por la pauta: Rs y n están acoplados (mal condicionado). "
             "Mitigación: bounds + inicialización analítica. Evidencia: nI=1.20 en ambas tecnologías."),
    'perdidas_pr': (55, "El gráfico clave: la pendiente térmica de m-Si es el doble de empinada. "
                    "El calor castiga justo al mediodía, cuando más energía hay en juego."),
    'veredicto_eco': (65, "Veredicto del track validado (recurso emulado, R²=0.99): HIT 84.18 % vs m-Si "
                      "81.61 % → +2.57 puntos. La caída frente al deck previo es el ~3 % óptico que ahora SÍ "
                      "penalizamos. La prima CAPEX de HIT se justifica."),
    'recurso_real': (55, "Con el recurso REAL de Atacama (PVGIS TMY): yields ~2,360-2,390 kWh/kWp (2.3× el "
                     "emulado). HIT gana +1.17 pts de PR — menos que en Florida porque el aire de altura es "
                     "frío (Tc ~58-62 °C) — pero son +3,294 MWh y +USD 148k al año en 100 MWp."),
    'cumplimiento': (40, "Cierre de consultoría: cada requisito de la pauta tiene método, lámina de "
                     "evidencia y resultado. Nada quedó sin responder."),
    'limitaciones': (60, "Transparencia: resolvimos IAM/espectral, albedo 0.20, doble cobertura y recurso "
                     "real. Permanece: validación eléctrica con datos de Florida, α de literatura, sin soiling ni BOS."),
    'conclusiones': (60, "HIT gana por física (menor coeficiente térmico) con metodología validada e "
                     "IAM/espectral aplicados. Futuro: validación in-situ, soiling, BOS y doble diodo."),
    'referencias': (10, "Referencias disponibles. Gracias — abrimos la ronda de preguntas."),
}


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    acum = 0.0
    for key in ORDER:
        CURRENT_KEY[0] = key
        slide = prs.slides.add_slide(blank)
        apply_bg(slide)
        BUILDERS[key](slide, REF[key])
        add_footer(slide, REF[key])
        if key in NOTES:
            secs, msg = NOTES[key]
            ini = acum
            acum += secs
            marca = f"[min {ini/60:.1f} → {acum/60:.1f} | {secs} s]"
            slide.notes_slide.notes_text_frame.text = f"{marca}\n{msg}"
        elif key.startswith('anexo'):
            slide.notes_slide.notes_text_frame.text = (
                "ANEXO — no se presenta en los 20 min; respaldo para la ronda de preguntas.")
        print(f"  [{REF[key]:>2}/{TOTAL}] {key}")
    print(f"\nGuion principal: {acum/60:.1f} min para {len(NOTES)} láminas presentadas.")
    out = 'output/Presentacion_Final_ELI556_Atacama_POA_Tarea2_revisado.pptx'
    try:
        prs.save(out)
    except PermissionError:
        out = out.replace('.pptx', '_v2.pptx')
        prs.save(out)
    print(f"\nPresentación guardada en: {out}")


if __name__ == '__main__':
    create_presentation()
